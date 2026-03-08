#!/usr/bin/env python3
"""
Extract structured content from Hebrew lesson PPTX files.

Usage:
    python extract_pptx.py <input.pptx> <output_dir>

Outputs:
    <output_dir>/extracted.md   — structured markdown with all content
    <output_dir>/images/        — extracted images from slides
"""

import sys
import os
import re
import base64
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx --break-system-packages")
    sys.exit(1)


# ---------------------------------------------------------------------------
# Helpers used by classify_slide
# ---------------------------------------------------------------------------

def _has_underscores(text):
    """Check if text contains fill-in blanks (2+ consecutive underscores)."""
    return bool(re.search(r'_{2,}', text))


def _count_empty_table_cells(table_data):
    """Count empty cells across all tables (excluding header row)."""
    total = 0
    for table in table_data:
        for row in table[1:]:  # skip header
            for cell in row:
                if cell.strip() == "":
                    total += 1
    return total


def _has_multiple_choice(text):
    """Detect multiple-choice options in parentheses: (opt1, opt2, opt3)."""
    # Hebrew options in parentheses with commas — at least 2 options
    matches = re.findall(r'\([^)]*[\u0590-\u05FF][^)]*,\s*[^)]*[\u0590-\u05FF][^)]*\)', text)
    return len(matches) >= 2


def _count_russian_sentences(text):
    """Count standalone Russian sentences (not short grammar labels)."""
    # Numbered Russian lines
    numbered = re.findall(r'^\d+[\.\)]\s*[А-Яа-яЁё]', text, re.MULTILINE)
    return len(numbered)


def _is_numbered_grammar_list(text):
    """Check if numbered Russian items look like grammar headings (short + colon)."""
    items = re.findall(r'^\d+[\.\)]\s*([А-Яа-яЁё][^\n]{0,60})', text, re.MULTILINE)
    if not items:
        return False
    # If most items end with ':' or are very short — grammar headings, not translations
    colon_count = sum(1 for item in items if ':' in item[:50])
    short_count = sum(1 for item in items if len(item.strip()) < 30)
    return colon_count >= len(items) * 0.5 or (short_count >= len(items) * 0.7 and len(items) <= 5)


def _count_unnumbered_russian_sentences(text):
    """Count lines that are standalone Russian sentences (no number prefix)."""
    lines = text.split('\n')
    count = 0
    for line in lines:
        line = line.strip()
        if not line:
            continue
        ru_chars = len(re.findall(r'[А-Яа-яЁё]', line))
        he_chars = len(re.findall(r'[\u0590-\u05FF]', line))
        # Predominantly Russian line, sentence-length
        if ru_chars > 20 and ru_chars > he_chars * 2 and len(line) > 30:
            count += 1
    return count


# ---------------------------------------------------------------------------
# Main classifier
# ---------------------------------------------------------------------------

def classify_slide(texts, has_table, has_image, table_data):
    """
    Classify a slide into a category. Categories:
    - grammar: Russian explanatory text with Hebrew examples
    - vocabulary: table with מילים חדשות header
    - conjugation: conjugation table (שורש, עבר, הווה, etc.)
    - verb_summary: summary table of verbs by gzarot (שלמים, פ״נ, etc.)
    - shem_peula: verbal noun table (שם פעולה)
    - numerals: numeral reference tables (שם מספר)
    - prepositions: preposition declension tables (מילת יחס + suffixes)
    - exercise_fill: fill-in-the-blank exercises
    - exercise_translate: translation exercises
    - exercise_table: table-based exercises (fill table cells)
    - exercise_choose: choose the correct word from options
    - exercise_oged: אוגד exercises (specifically)
    - exercise_morphology: determine root/binyan/gzara from a verb form
    - exercise_transform: rewrite sentence in different tense (no blanks)
    - exercise_writing: write a composition (חיבור)
    - exercise_questions: answer questions about a text
    - reading: reading text (long Hebrew text, dialog, story)
    - image_task: slide with primarily an image (spot-the-difference, etc.)
    - joke: jokes / humorous texts
    - other: unclassified
    """
    all_text = "\n".join(texts).strip()
    all_text_lower = all_text.lower()
    text_len_stripped = len(all_text.replace(" ", "").replace("\n", ""))

    russian_chars = len(re.findall(r'[А-Яа-яЁё]', all_text))
    hebrew_chars = len(re.findall(r'[\u0590-\u05FF]', all_text))
    has_blanks = _has_underscores(all_text)

    # ---- Grammar keyword lists (used in multiple phases) ----
    grammar_keywords_ru = [
        "связк", "подлежащ", "сказуем", "именн", "копул", "огэд",
        "Правил", "спряжени", "глагол", "породы", "породе", "корн",
        "префикс", "огласов", "биньян", "гзар", "каузатив",
        "Итак", "Дело в том", "Приведем", "означает",
        "Признак", "Обратите вниман", "Рассмотрим",
        "Продолжим", "Вспомним", "Именно это",
        "модел", "образова", "конструкц", "сопряжён",
        "Сравнител", "Похож", "используется", "используются",
        "прилагательн", "существительн", "числительн",
        "ассимиляц", "артикл", "инфинитив",
        "В иврите", "В русском языке",
    ]
    grammar_keywords_he = [
        "אוגד", "דגוא", "משפט שמני", "ינמש טפשמ",
        "בניין", "ןיינב", "גזרת", "תרזג", "הפעיל", "ליעפה",
        "התפעל", "לעפתה", "פיעל", "לעיפ", "נפעל", "לעפנ",
    ]

    # ---- Gather table metadata ----
    all_headers = []
    all_cells_text = []
    total_empty_cells = 0
    if has_table and table_data:
        for table in table_data:
            if table and len(table) > 0:
                all_headers.extend([cell.strip() for cell in table[0]])
                for row in table:
                    all_cells_text.extend([cell.strip() for cell in row])
        total_empty_cells = _count_empty_table_cells(table_data)
    headers_joined = " ".join(all_headers)
    cells_joined = " ".join(all_cells_text)

    # Also check for blanks inside table cells
    table_has_blanks = any(_has_underscores(c) for c in all_cells_text)

    # ==================================================================
    # PHASE 1: Image-only slides (very little text, primarily image)
    # BUT: do NOT classify as image_task if there's a meaningful table
    # ==================================================================
    if has_image and text_len_stripped < 30 and not has_table:
        return "image_task"

    # ==================================================================
    # PHASE 2: Instruction-keyword based exercise detection (EARLY)
    # These keywords reliably signal exercises regardless of text length
    # ==================================================================

    # --- Exercise: writing / composition ---
    if "תכתבו חיבור" in all_text or "תכתוב חיבור" in all_text or "רוביח ובתכת" in all_text:
        return "exercise_writing"

    # --- Exercise: questions about a text ---
    if ("תשאלו" in all_text or "ענה/עני על השאלות" in all_text or
            "ענו על השאלות" in all_text or "תולאשה לע" in all_text):
        # Only if has numbered questions or bullet points
        if re.search(r'[•\d]', all_text):
            return "exercise_questions"

    # --- Exercise: morphology (determine root/binyan/gzara) ---
    if "תעשו לפי הדוגמה" in all_text or "המגודה יפל ושעת" in all_text:
        # Check for pattern: verb form + blanks for root/binyan/gzara
        if has_blanks:
            return "exercise_morphology"

    # --- Exercise: translation (both תתרגמו and תרגמו) ---
    if re.search(r'ת?תרגמו', all_text) or "ומגרתת" in all_text or "ומגרת" in all_text:
        return "exercise_translate"

    # --- Exercise: oged (ONLY when אוגד is explicitly mentioned) ---
    if ("אוגד" in all_text or "דגוא" in all_text) and has_blanks:
        return "exercise_oged"

    # --- Exercise: choose correct word (תבחרו / multiple choice in parens) ---
    if "תבחרו" in all_text or "ורחבת" in all_text:
        return "exercise_choose"
    if "הנוכנה הלימב" in all_text or "הפועל הנכון" in all_text or "הצורה הנכונה" in all_text:
        return "exercise_choose"
    if _has_multiple_choice(all_text) and has_blanks:
        return "exercise_choose"

    # ==================================================================
    # PHASE 3: Table-based classification
    # ==================================================================
    if has_table and table_data:

        # --- Numerals table (שם מספר) ---
        if "שם מספר" in all_text or "רפסמ םש" in all_text:
            return "numerals"
        # Tables where first column is numbers (10, 20, 100, 1000...)
        first_col_nums = 0
        for table in table_data:
            for row in table[1:]:
                if row and re.match(r'^\d+$', row[0].strip()):
                    first_col_nums += 1
        if first_col_nums >= 5:
            # Likely a numerals reference table
            # But confirm it's not something else
            if not any(kw in headers_joined for kw in ["שורש", "שם פועל", "פעולה"]):
                return "numerals"

        # --- Prepositions table (מילת היחס + declension) ---
        if "מילת היחס" in all_text or "סחיה תלימ" in all_text:
            return "prepositions"
        # Russian gender/number headers → preposition declension
        if any(kw in headers_joined for kw in ["Ед. число", "Мн. число", "Жен. род", "Муж. род"]):
            return "prepositions"

        # --- Conjugation table ---
        if any(kw in headers_joined for kw in ["שורש", "יחיד זכר", "רבים זכר", "יחיד נקבה"]):
            return "conjugation"

        # --- Tense-based table (עבר, הווה, עתיד) ---
        if any(kw in headers_joined for kw in ["עבר", "הווה", "עתיד"]):
            if ("תמלאו" in all_text or "ואלמת" in all_text or
                    table_has_blanks or total_empty_cells > 3):
                return "exercise_table"
            return "conjugation"

        # --- Shem peula table ---
        if any(kw in headers_joined for kw in ["שם פעולה", "שם פועלה", "פעולה"]):
            if "תרגום" in headers_joined or "שם" in headers_joined:
                return "shem_peula"
        if "שם הפועלה" in all_text or "הלעופה םש" in all_text:
            return "shem_peula"

        # --- Verb summary by gzarot ---
        if any(kw in headers_joined for kw in ["שלמים", "גרזה", "גזרה"]):
            return "verb_summary"
        # Check cells for gzara markers
        if any(kw in cells_joined for kw in ["שלמים", "פ״נ", "פ״י", "ע״ו"]):
            if any(kw in all_text for kw in ["גזרות", "תורזג", "בניין", "ןיינב"]):
                return "verb_summary"

        # --- Vocabulary table ---
        vocab_markers = ["שם תואר", "פעלים", "מיליות", "ביטויים"]
        if any(kw in headers_joined for kw in vocab_markers):
            return "vocabulary"
        if "שם פועל" in headers_joined and any(kw in headers_joined for kw in ["תואר", "פעלים"]):
            return "vocabulary"

        # --- Exercise table (by empty cells or instruction keywords) ---
        if "תמלאו" in all_text or "ואלמת" in all_text:
            return "exercise_table"
        # Pronoun-header tables with empty cells (אני/אתה/הוא pattern)
        pronoun_headers = ["אני", "אתה", "את", "הוא", "היא", "אנחנו"]
        if sum(1 for h in all_headers if h in pronoun_headers) >= 3 and total_empty_cells > 3:
            return "exercise_table"
        # שם פועל header with many empty cells = exercise table
        if "שם פועל" in headers_joined and total_empty_cells > 3:
            return "exercise_table"
        # Generic: table with many empty cells + exercise instruction
        if total_empty_cells > 3 and any(kw in all_text for kw in ["תמלאו", "תשלימו", "תאמרו"]):
            return "exercise_table"
        # Table with many empty cells and no known reference pattern
        if total_empty_cells > 5:
            return "exercise_table"

        # --- Grammar table with only Russian text (e.g. "Никогда" table) ---
        ru_in_cells = len(re.findall(r'[А-Яа-яЁё]', cells_joined))
        he_in_cells = len(re.findall(r'[\u0590-\u05FF]', cells_joined))
        if ru_in_cells > 30 and (hebrew_chars > 10 or he_in_cells > 10):
            return "grammar"

    # ==================================================================
    # PHASE 4: Text-based exercise detection (blanks in text)
    # ==================================================================
    if has_blanks or table_has_blanks:
        if has_table:
            return "exercise_table"

        if "תשלימו" in all_text or "ומילשת" in all_text:
            return "exercise_fill"
        if "תעשו משפטים" in all_text or "טפשמ ושעת" in all_text:
            return "exercise_fill"
        if "השלימו" in all_text or "ומילשה" in all_text:
            return "exercise_fill"
        if "תכתבו" in all_text or "ובתכת" in all_text:
            return "exercise_fill"
        if "שבצו" in all_text:
            return "exercise_fill"

        # הוסיפו WITHOUT אוגד = regular fill exercise
        if "הוסיפו" in all_text or "ופיסוה" in all_text:
            return "exercise_fill"

        return "exercise_fill"

    # ==================================================================
    # PHASE 5: Translation exercise detection (no blanks)
    # ==================================================================

    # Numbered Russian sentences → translation (but filter grammar headings)
    num_ru_sentences = _count_russian_sentences(all_text)
    if num_ru_sentences >= 3:
        if not _is_numbered_grammar_list(all_text):
            # Confirm no strong grammar signals
            grammar_exclude = ["Правил", "Дело в том", "Итак", "К данному",
                               "Приведем", "Обратите вниман", "Рассмотрим"]
            if not any(kw in all_text for kw in grammar_exclude):
                return "exercise_translate"

    # Unnumbered Russian sentences (≥5 standalone Russian lines)
    if _count_unnumbered_russian_sentences(all_text) >= 5:
        # Check for translation instruction nearby
        if re.search(r'תרגמו|ומגרת', all_text):
            return "exercise_translate"
        # Many Russian sentences with minimal Hebrew = likely translation
        # BUT: skip if grammar keywords are present (it's an explanation, not exercise)
        if russian_chars > 300 and hebrew_chars < 50:
            has_grammar_signal = (
                any(kw in all_text for kw in grammar_keywords_ru) or
                any(kw in all_text for kw in grammar_keywords_he)
            )
            if not has_grammar_signal:
                return "exercise_translate"

    # ==================================================================
    # PHASE 6: Exercise transform (rewrite in different tense, no blanks)
    # ==================================================================
    if "תשלימו" in all_text or "ומילשת" in all_text:
        # Has instruction but NO blanks — likely a transform exercise
        if num_ru_sentences == 0:  # Hebrew-only
            return "exercise_transform"
        # Has numbered Hebrew sentences to rewrite
        he_numbered = re.findall(r'^\d+[\.\)]\s*[\u0590-\u05FF]', all_text, re.MULTILINE)
        if len(he_numbered) >= 3:
            return "exercise_transform"

    # ==================================================================
    # PHASE 7: Grammar detection
    # ==================================================================
    if russian_chars > 50:
        if any(kw in all_text for kw in grammar_keywords_ru):
            return "grammar"
        if any(kw in all_text for kw in grammar_keywords_he):
            return "grammar"

    # ==================================================================
    # PHASE 8: Jokes
    # ==================================================================
    joke_markers = ["גולם", "חושם", "חלם", "לבהא", "חידה", "חידות",
                    "םלוג", "םשוח", "םלח"]
    if any(m in all_text for m in joke_markers):
        if len(all_text) > 200:
            return "joke"

    # ==================================================================
    # PHASE 9: Reading text (long Hebrew, dialog, story)
    # ==================================================================
    if len(all_text) > 500:
        # But first — check if it's actually grammar with lots of text
        if russian_chars > 200 and russian_chars > hebrew_chars * 0.3:
            return "grammar"

        # Dialog markers
        if all_text.count("-") >= 3 or all_text.count("—") >= 3:
            return "reading"
        # Primarily Hebrew
        total_alpha = len(re.findall(r'[\w]', all_text))
        if total_alpha > 0 and hebrew_chars / total_alpha > 0.5:
            return "reading"

    # ==================================================================
    # PHASE 10: Shem peula (standalone, not in table)
    # ==================================================================
    if "שם פעולה" in all_text or "הלעופ םש" in all_text:
        return "shem_peula"
    if "שם הפועלה" in all_text or "הלעופה םש" in all_text:
        return "shem_peula"

    # ==================================================================
    # PHASE 11: Fallbacks
    # ==================================================================

    # Grammar fallback: substantial Russian with some Hebrew
    if russian_chars > 50 and hebrew_chars > 10:
        # Re-check for numbered Russian → translate
        if num_ru_sentences >= 3 and not _is_numbered_grammar_list(all_text):
            return "exercise_translate"
        if russian_chars > 200:
            return "grammar"

    # Table without classification yet
    if has_table and table_data:
        for table in table_data:
            for row in table:
                for cell in row:
                    if _has_underscores(cell):
                        return "exercise_table"
        # Table with substantial Hebrew content
        if hebrew_chars > 50:
            return "grammar"

    # Moderate-length content
    if len(all_text) > 100:
        if russian_chars > hebrew_chars:
            return "grammar"
        return "reading"

    # Image with some text
    if has_image:
        return "image_task"

    return "other"


def extract_table(shape):
    """Extract table data as list of lists."""
    table = shape.table
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
    return rows


def extract_images(slide, slide_idx, output_dir):
    """Extract images from a slide, save to files, return list of paths."""
    images_dir = os.path.join(output_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    image_paths = []
    img_count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_count += 1
            image = shape.image
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            filename = f"slide_{slide_idx + 1}_img_{img_count}.{ext}"
            filepath = os.path.join(images_dir, filename)
            with open(filepath, "wb") as f:
                f.write(image.blob)
            image_paths.append(filepath)

    return image_paths


def extract_slide_content(slide, slide_idx, output_dir):
    """Extract all content from a single slide."""
    texts = []
    tables = []
    has_image = False
    image_paths = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_image = True

        if shape.has_text_frame:
            # Skip slide number and footer placeholders
            shape_name = shape.name if hasattr(shape, 'name') else ""
            is_placeholder = ("מספר שקופית" in shape_name or
                            "שקופית" in shape_name and "מציין" in shape_name or
                            "כותרת תחתונה" in shape_name)
            if is_placeholder:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)

        if shape.has_table:
            tables.append(extract_table(shape))

    if has_image:
        image_paths = extract_images(slide, slide_idx, output_dir)

    category = classify_slide(texts, bool(tables), has_image, tables)

    return {
        "slide_number": slide_idx + 1,
        "category": category,
        "texts": texts,
        "tables": tables,
        "has_image": has_image,
        "image_paths": image_paths,
    }


def table_to_markdown(table_data):
    """Convert table data to markdown table."""
    if not table_data:
        return ""

    lines = []
    # Header
    header = table_data[0]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")

    # Rows
    for row in table_data[1:]:
        # Pad row if shorter than header
        padded = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(padded[:len(header)]) + " |")

    return "\n".join(lines)


def format_slide_md(slide_data):
    """Format a single slide's data as markdown."""
    lines = []
    cat = slide_data["category"]
    num = slide_data["slide_number"]

    # Category label
    category_labels = {
        "grammar": "📖 Грамматика",
        "vocabulary": "📝 Новые слова",
        "conjugation": "📊 Таблица спряжения",
        "verb_summary": "📊 Сводная таблица глаголов",
        "shem_peula": "📊 Отглагольные существительные (שם פעולה)",
        "numerals": "🔢 Числительные (שם מספר)",
        "prepositions": "📖 Предлоги (מילות יחס)",
        "exercise_fill": "✏️ Упражнение: вставить слово",
        "exercise_translate": "✏️ Упражнение: перевод",
        "exercise_table": "✏️ Упражнение: заполнить таблицу",
        "exercise_choose": "✏️ Упражнение: выбрать слово",
        "exercise_oged": "✏️ Упражнение: אוגד",
        "exercise_morphology": "✏️ Упражнение: морфологический разбор",
        "exercise_transform": "✏️ Упражнение: трансформация",
        "exercise_writing": "✏️ Упражнение: сочинение (חיבור)",
        "exercise_questions": "✏️ Упражнение: вопросы к тексту",
        "reading": "📖 Текст для чтения",
        "image_task": "🖼️ Задание по картинке",
        "joke": "😄 Анекдоты / юмор",
        "other": "📄 Прочее",
    }

    label = category_labels.get(cat, cat)
    lines.append(f"## Слайд {num} — {label}")
    lines.append("")

    # Image placeholder
    if slide_data["has_image"] and slide_data["image_paths"]:
        for img_path in slide_data["image_paths"]:
            rel_path = os.path.basename(img_path)
            lines.append(f"![Изображение](images/{rel_path})")
            lines.append("")
        if cat == "image_task":
            lines.append("<!-- CLAUDE: Опиши изображение и найди различия между двумя частями картинки -->")
            lines.append("")

    # Text content
    if slide_data["texts"]:
        for text in slide_data["texts"]:
            lines.append(text)
            lines.append("")

    # Tables
    if slide_data["tables"]:
        for table in slide_data["tables"]:
            lines.append(table_to_markdown(table))
            lines.append("")

    lines.append("---")
    lines.append("")

    return "\n".join(lines)


def extract_lesson_info(slides_data, filename=None):
    """Try to determine lesson number and part from slide content or filename."""
    # First try from slide content
    for slide in slides_data:
        for text in slide["texts"]:
            match = re.search(r'שיעור\s+(\d+)\s+חלק\s+(\d+)', text)
            if match:
                return int(match.group(1)), int(match.group(2))
            match = re.search(r'(\d+)\s+קלח\s+(\d+)\s+רועיש', text)
            if match:
                return int(match.group(2)), int(match.group(1))

    # Fallback: try filename
    if filename:
        match = re.search(r'שיעור[_\s]+(\d+)[_\s]+חלק[_\s]+(\d+)', filename)
        if match:
            return int(match.group(1)), int(match.group(2))
        match = re.search(r'(\d+)[_\s]+חלק[_\s]+(\d+)', filename)
        if match:
            return int(match.group(1)), int(match.group(2))

    return None, None


def main():
    if len(sys.argv) < 3:
        print(f"Usage: {sys.argv[0]} <input.pptx> <output_dir>")
        sys.exit(1)

    input_path = sys.argv[1]
    output_dir = sys.argv[2]

    if not os.path.exists(input_path):
        print(f"ERROR: File not found: {input_path}")
        sys.exit(1)

    os.makedirs(output_dir, exist_ok=True)

    print(f"Opening: {input_path}")
    prs = Presentation(input_path)

    slides_data = []
    for idx, slide in enumerate(prs.slides):
        print(f"  Processing slide {idx + 1}...")
        data = extract_slide_content(slide, idx, output_dir)
        slides_data.append(data)
        print(f"    → {data['category']} ({len(data['texts'])} text blocks, {len(data['tables'])} tables, images: {data['has_image']})")

    # Determine lesson info
    filename = os.path.basename(input_path)
    lesson_num, part_num = extract_lesson_info(slides_data, filename)

    # Build markdown
    md_lines = []

    if lesson_num and part_num:
        md_lines.append(f"# Извлечение из урока {lesson_num}, часть {part_num}")
        md_lines.append(f"## שיעור {lesson_num} חלק {part_num}")
    else:
        md_lines.append("# Извлечение из урока")

    md_lines.append("")
    md_lines.append(f"Всего слайдов: {len(slides_data)}")
    md_lines.append("")

    # Summary of slide categories
    md_lines.append("### Структура презентации")
    md_lines.append("")
    for slide in slides_data:
        cat = slide["category"]
        label_map = {
            "grammar": "грамматика",
            "vocabulary": "словарь",
            "conjugation": "спряжение",
            "verb_summary": "сводка глаголов",
            "shem_peula": "שם פעולה",
            "numerals": "числительные",
            "prepositions": "предлоги",
            "exercise_fill": "упражнение (вставить)",
            "exercise_translate": "упражнение (перевод)",
            "exercise_table": "упражнение (таблица)",
            "exercise_choose": "упражнение (выбор)",
            "exercise_oged": "упражнение (אוגד)",
            "exercise_morphology": "упражнение (морфология)",
            "exercise_transform": "упражнение (трансформация)",
            "exercise_writing": "упражнение (сочинение)",
            "exercise_questions": "упражнение (вопросы)",
            "reading": "текст",
            "image_task": "картинка",
            "joke": "анекдоты",
            "other": "прочее",
        }
        md_lines.append(f"- Слайд {slide['slide_number']}: {label_map.get(cat, cat)}")
    md_lines.append("")
    md_lines.append("---")
    md_lines.append("")

    # All slides
    for slide in slides_data:
        md_lines.append(format_slide_md(slide))

    # Build output filename based on lesson/part numbers
    if lesson_num is not None and part_num is not None:
        base_name = f"урок_{lesson_num}_часть_{part_num}_extracted"
    else:
        base_name = "extracted"

    # Write output
    md_path = os.path.join(output_dir, f"{base_name}.md")
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))

    print(f"\nDone!")
    print(f"  Markdown: {md_path}")
    if any(s["has_image"] for s in slides_data):
        print(f"  Images:   {os.path.join(output_dir, 'images/')}")

    # Print category summary
    from collections import Counter
    cats = Counter(s["category"] for s in slides_data)
    print(f"\nCategory summary:")
    for cat, count in cats.most_common():
        print(f"  {cat}: {count}")


if __name__ == "__main__":
    main()
